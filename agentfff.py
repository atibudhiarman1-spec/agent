import argparse, os, tempfile, uuid, json
from typing import TypedDict, List, Dict, Optional, Literal

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

from pydantic import BaseModel
from dotenv import load_dotenv
from langgraph.graph import StateGraph, END
from langchain_google_genai import ChatGoogleGenerativeAI

load_dotenv()
API_KEY = os.getenv("GOOGLE_API_KEY") or os.getenv("GEMINI_API_KEY")
if not API_KEY:
    raise RuntimeError("Set GOOGLE_API_KEY (or GEMINI_API_KEY) in your .env")
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")

def llm():
    return ChatGoogleGenerativeAI(model=GEMINI_MODEL, api_key=API_KEY, temperature=0)

class State(TypedDict, total=False):
    excel_path: str
    ppt_template_path: str
    mapping_path: Optional[str]
    mapping: Dict
    df: pd.DataFrame
    tempdir: str
    chart_plan: dict
    charts: List[Dict[str, object]]
    insights: List[str]
    ppt_out_path: str
    log: List[str]

class ChartSpec(BaseModel):
    id: str
    type: Literal["line", "bar"]
    x: str
    y: str
    agg: Literal["sum", "mean", "none"] = "sum"
    top_k: Optional[int] = None

class ChartPlan(BaseModel):
    charts: List[ChartSpec]

class InsightPayload(BaseModel):
    summary: str
    bullets: List[str]

def _df_profile(df: pd.DataFrame) -> dict:
    prof = {"columns": []}
    for col in df.columns:
        series = df[col]
        entry = {"name": str(col), "dtype": str(series.dtype)}
        if np.issubdtype(series.dtype, np.number):
            s = pd.to_numeric(series, errors="coerce")
            entry["stats"] = {
                "count": int(s.count()),
                "sum": float(np.nansum(s)),
                "mean": float(np.nanmean(s)) if s.count() else 0.0,
                "min": float(np.nanmin(s)) if s.count() else 0.0,
                "max": float(np.nanmax(s)) if s.count() else 0.0,
            }
        elif np.issubdtype(series.dtype, np.datetime64):
            s = pd.to_datetime(series, errors="coerce").dropna()
            entry["time_coverage"] = {
                "min": s.min().isoformat() if len(s) else None,
                "max": s.max().isoformat() if len(s) else None,
            }
        else:
            vc = series.astype(str).value_counts().head(10)
            entry["top_values"] = [{"v": k, "n": int(v)} for k, v in vc.items()]
        prof["columns"].append(entry)
    return prof

def _load_mapping_file(path: Optional[str]) -> Dict:
    """
    Supports YAML (.yaml/.yml) and JSON (.json).

    Expected keys (all OPTIONAL):
      charts: {<chart_id>: <shape_name>, ...}
      insights: <shape_name>
      requirements:
        <chart_id>: <"line"|"bar">
        or
        <chart_id>:
          type: "line"|"bar"
          x: "<column name>"
          y: "<column name>"
          agg: "sum"|"mean"|"none"
          top_k: <int>
      place_unmapped: true|false
    """
    if not path:
        return {}
    if not os.path.exists(path):
        raise FileNotFoundError(f"Mapping file not found: {path}")
    ext = os.path.splitext(path)[1].lower()
    if ext in (".yml", ".yaml"):
        try:
            import yaml  # type: ignore
        except Exception as e:
            raise RuntimeError("pyyaml not installed. Run: pip install pyyaml") from e
        with open(path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f) or {}
    elif ext == ".json":
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f) or {}
    else:
        raise ValueError("Mapping file must be .yaml/.yml or .json")
    data.setdefault("charts", {})
    data.setdefault("requirements", {})
    return data

def _pick_numeric(df: pd.DataFrame) -> Optional[str]:
    nums = df.select_dtypes(include="number").columns.tolist()
    if not nums:
        return None
    return max(nums, key=lambda c: pd.to_numeric(df[c], errors="coerce").fillna(0.0).var())

def _pick_datetime(df: pd.DataFrame) -> Optional[str]:
    candidates = [c for c in df.columns if np.issubdtype(df[c].dtype, np.datetime64)]
    return candidates[0] if candidates else None

def _pick_categorical(df: pd.DataFrame) -> Optional[str]:
    cats = df.select_dtypes(include=["object", "category"]).columns.tolist()
    return cats[0] if cats else None

def _coerce_series(vals) -> List[float]:
    s = pd.to_numeric(pd.Series(vals), errors="coerce").fillna(0.0)
    return [float(v) for v in s.tolist()]

def _xl_chart_type(kind: str) -> XL_CHART_TYPE:
    return XL_CHART_TYPE.LINE_MARKERS if kind == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED

def _safe_chart_data(payload: Dict[str, object]) -> ChartData:
    """Build ChartData safely: non-empty, aligned categories/values."""
    cats = [str(x) for x in (payload.get("categories") or [])]
    vals = [float(v) for v in (payload.get("values") or [])]
    if not cats or not vals:
        cats, vals = ["N/A"], [0.0]
    n = min(len(cats), len(vals))
    cats, vals = cats[:n], vals[:n]
    cd = ChartData()
    cd.categories = cats
    cd.add_series(str(payload.get("series_name", "Series 1")), vals)
    return cd

def validate_inputs(s: State) -> State:
    log = s.get("log", [])
    if not s.get("excel_path") or not os.path.exists(s["excel_path"]):
        raise FileNotFoundError(f"Excel file not found: {s.get('excel_path')}")
    if not s.get("ppt_template_path") or not os.path.exists(s["ppt_template_path"]):
        raise FileNotFoundError(f"PPT template not found: {s.get('ppt_template_path')}")
    tempdir = tempfile.mkdtemp(prefix="lg_ppt_")
    log.append(f"[validate] tempdir={tempdir}")
    return {"tempdir": tempdir, "log": log}

def load_mapping(s: State) -> State:
    log = s.get("log", [])
    mapping = _load_mapping_file(s.get("mapping_path"))
    charts = mapping.get("charts", {})
    insights = mapping.get("insights")
    reqs = mapping.get("requirements", {})
    log.append(f"[mapping] charts={len(charts)} insights={'set' if insights else 'none'} reqs={len(reqs)}")
    return {"mapping": mapping, "log": log}

def load_excel(s: State) -> State:
    log = s.get("log", [])
    df = pd.read_excel(s["excel_path"], engine="openpyxl", sheet_name=0)
    for col in df.columns:
        if "date" in str(col).lower():
            df[col] = pd.to_datetime(df[col], errors="coerce")
    log.append(f"[excel] shape={df.shape}")
    return {"df": df, "log": log}

def plan_charts_with_gemini(s: State) -> State:
    """Ask Gemini for a tiny chart plan (1–2 charts) using a schema-only summary."""
    log = s.get("log", [])
    schema = _df_profile(s["df"])
    prompt = (
        "Propose 1-2 charts from this dataframe schema.\n"
        "- LINE: datetime x + numeric y (aggregate daily)\n"
        "- BAR: categorical x + numeric y (top_k <= 10)\n"
        "- Return JSON matching the response schema only.\n\n"
        f"SCHEMA:\n{json.dumps(schema)}"
    )
    chart_planner = llm().with_structured_output(ChartPlan)
    try:
        plan = chart_planner.invoke(prompt)
        plan_dict = plan.model_dump()
        log.append(f"[plan:charts] LLM specs={len(plan_dict.get('charts', []))}")
        return {"chart_plan": plan_dict, "log": log}
    except Exception as e:
        df = s["df"]
        charts = []
        dt = _pick_datetime(df)
        num = _pick_numeric(df)
        if dt and num:
            charts.append(ChartSpec(id="ts", type="line", x=dt, y=num).model_dump())
        elif num:
            cat = _pick_categorical(df)
            if cat:
                charts.append(ChartSpec(id="bar1", type="bar", x=cat, y=num, top_k=10).model_dump())
        log.append(f"[plan:charts] fallback ({type(e).__name__})")
        return {"chart_plan": {"charts": charts}, "log": log}

def enforce_requirements_overlay(s: State) -> State:
    """
    Overlay OPTIONAL requirements on top of the LLM plan.
    - If requirement is string ("line"/"bar"), enforce type and synthesize columns.
    - If requirement is object ({type,x,y,agg,top_k}), enforce exactly (with validation).
    """
    log = s.get("log", [])
    df = s["df"]
    plan = ChartPlan(**s["chart_plan"])
    specs_by_id: Dict[str, ChartSpec] = {c.id: c for c in plan.charts}
    reqs: Dict[str, object] = (s.get("mapping") or {}).get("requirements", {}) or {}
    if not reqs:
        return {"chart_plan": plan.model_dump(), "log": log}
    for cid, req in reqs.items():
        if isinstance(req, str):
            req_type = req.strip().lower()
            if req_type not in ("line", "bar"):
                log.append(f"[require] '{cid}' invalid type='{req}' → ignored")
                continue
            if req_type == "line":
                x = _pick_datetime(df)
                y = _pick_numeric(df)
                if not x or not y:
                    y = _pick_numeric(df) or list(df.select_dtypes(include="number").columns[:1] or ["value"])[0]
                    x = x or str(next(iter(df.columns)))
                    log.append(f"[require] '{cid}' line: used heuristics x='{x}', y='{y}'")
                specs_by_id[cid] = ChartSpec(id=cid, type="line", x=x, y=y)
            else:
                x = _pick_categorical(df) or str(next(iter(df.columns)))
                y = _pick_numeric(df) or x
                specs_by_id[cid] = ChartSpec(id=cid, type="bar", x=x, y=y, top_k=10)
                log.append(f"[require] '{cid}' bar: used heuristics x='{x}', y='{y}'")
        elif isinstance(req, dict):
            t = str(req.get("type", "")).lower()
            x = req.get("x")
            y = req.get("y")
            agg = req.get("agg", "sum")
            top_k = req.get("top_k", None)
            if t not in ("line", "bar"):
                log.append(f"[require] '{cid}' invalid dict.type='{t}' → ignored")
                continue
            missing = [col for col in (x, y) if col is not None and col not in df.columns]
            if missing:
                log.append(f"[require] '{cid}' missing columns {missing} → will heuristically fix")
            if t == "line":
                x = x if (x in df.columns and np.issubdtype(df[x].dtype, np.datetime64)) else (_pick_datetime(df) or x or str(next(iter(df.columns))))
                y = y if (y in df.columns) else (_pick_numeric(df) or y or str(next(iter(df.columns))))
                specs_by_id[cid] = ChartSpec(id=cid, type="line", x=str(x), y=str(y), agg=agg)
            else:
                x = x if (x in df.columns) else (_pick_categorical(df) or x or str(next(iter(df.columns))))
                y = y if (y in df.columns) else (_pick_numeric(df) or y or str(next(iter(df.columns))))
                specs_by_id[cid] = ChartSpec(id=cid, type="bar", x=str(x), y=str(y), top_k=top_k, agg=agg)
            log.append(f"[require] '{cid}' enforced: type={t}, x='{x}', y='{y}'")
        else:
            log.append(f"[require] '{cid}' invalid requirement format → ignored")
    updated = ChartPlan(charts=list(specs_by_id.values())).model_dump()
    return {"chart_plan": updated, "log": log}

def render_charts_prepare_payloads(s: State) -> State:
    """
    Prepare chart payloads for native PPT (categories/values) and also write a PNG fallback.
    """
    df = s["df"].copy()
    tempdir = s["tempdir"]
    plan = ChartPlan(**s["chart_plan"])
    out: List[Dict[str, object]] = []
    log = s.get("log", [])
    for spec in plan.charts:
        if spec.type == "line":
            x = df[spec.x] if spec.x in df.columns else pd.Series(range(len(df)))
            y = pd.to_numeric(df[spec.y], errors="coerce") if spec.y in df.columns else pd.Series(dtype=float)
            if np.issubdtype(x.dtype, np.datetime64):
                dfx = pd.DataFrame({"x": pd.to_datetime(x, errors="coerce"), "y": y}).dropna()
                if spec.agg == "mean":
                    dfx = dfx.groupby(dfx["x"].dt.date)["y"].mean().reset_index()
                else:
                    dfx = dfx.groupby(dfx["x"].dt.date)["y"].sum().reset_index()
                categories = [d.isoformat() for d in dfx["x"]]
                values = _coerce_series(dfx["y"])
                title = f"{spec.y} over time"
            else:
                y = pd.to_numeric(y, errors="coerce").fillna(0.0)
                categories = [str(i) for i in range(len(y))]
                values = _coerce_series(y)
                title = f"{spec.y} vs row"
            fig = plt.figure()
            plt.plot(range(len(values)), values)
            plt.title(title)
            plt.tight_layout()
            img = os.path.join(tempdir, f"{uuid.uuid4().hex}_line.png")
            fig.savefig(img, dpi=200)
            plt.close(fig)
            out.append({
                "id": spec.id, "type": "line", "title": title,
                "categories": categories, "values": values,
                "series_name": spec.y, "img_path": img
            })
        elif spec.type == "bar":
            x = df[spec.x].astype(str) if spec.x in df.columns else pd.Series([], dtype=str)
            y = pd.to_numeric(df[spec.y], errors="coerce") if spec.y in df.columns else pd.Series([], dtype=float)
            dfx = pd.DataFrame({"x": x, "y": y}).dropna()
            grp = dfx.groupby("x")["y"]
            if spec.agg == "mean":
                grp = grp.mean()
            else:
                grp = grp.sum()
            grp = grp.sort_values(ascending=False)
            if spec.top_k:
                grp = grp.head(spec.top_k)
            dfx = grp.reset_index()
            categories = dfx["x"].astype(str).tolist()
            values = _coerce_series(dfx["y"])
            title = f"Top {spec.x} by {spec.y}"
            fig = plt.figure()
            plt.bar(categories, values)
            plt.xticks(rotation=45, ha="right")
            plt.title(title)
            plt.tight_layout()
            img = os.path.join(tempdir, f"{uuid.uuid4().hex}_bar.png")
            fig.savefig(img, dpi=200)
            plt.close(fig)
            out.append({
                "id": spec.id, "type": "bar", "title": title,
                "categories": categories, "values": values,
                "series_name": spec.y, "img_path": img
            })
    log.append(f"[render] prepared {len(out)} chart(s) (native-ready, PNG fallback)")
    return {"charts": out, "log": log}

def generate_insights_gemini(s: State) -> State:
    df = s["df"]
    charts = s.get("charts", [])
    log = s.get("log", [])
    num_cols = df.select_dtypes(include="number").columns.tolist()
    aggregates = {}
    for c in num_cols[:6]:
        ser = pd.to_numeric(df[c], errors="coerce")
        aggregates[c] = {"sum": float(ser.sum()), "mean": float(ser.mean())}
    prompt = (
        "Produce a concise, business-friendly executive summary (1-2 sentences) plus 3-5 bullet insights. Insights should be standalone. Do not write things like available in chart."
        "Use ONLY the aggregates and chart titles provided. Include concrete numbers, deltas, or rankings when possible."
        f"\nAGGREGATES: {json.dumps(aggregates)}"
        f"\nCHARTS: {json.dumps([{'id':c['id'],'title':c['title']} for c in charts])}"
    )
    try:
        payload = llm().with_structured_output(InsightPayload).invoke(prompt)
        insights = [payload.summary, *payload.bullets]
    except Exception:
        insights = ["Insights unavailable due to a model error.", "This slide is a placeholder."]
    log.append(f"[insights] {len(insights)} line(s)")
    return {"insights": insights, "log": log}

def _find_shape_anywhere(prs: Presentation, name: str):
    for slide in prs.slides:
        for shp in slide.shapes:
            if getattr(shp, "name", None) == name:
                return slide, shp
    return None, None

def bind_ppt_native_with_mapping_and_fallback(s: State) -> State:
    prs = Presentation(s["ppt_template_path"])
    log = s.get("log", [])
    title_layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
    mapping = s.get("mapping") or {}
    mapped = (mapping.get("charts") or {})
    insights_target = mapping.get("insights") or "TEXT_insights"
    if "place_unmapped" in mapping:
        place_unmapped = bool(mapping["place_unmapped"])
    else:
        place_unmapped = False if mapped else True
    produced = list(s.get("charts", []))
    by_id = {c["id"]: c for c in produced}
    used: set = set()

    def add_native_chart_to_slide(slide, payload: Dict[str, object], left, top, width, height):
        try:
            chart_type = _xl_chart_type(str(payload["type"]))
            chart_data = _safe_chart_data(payload)
            graphic_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
            return getattr(graphic_frame, "chart", None) is not None
        except Exception as e:
            log.append(f"[native:add_chart:error] {type(e).__name__}: {e}")
            return False

    for cid, shape_name in mapped.items():
        payload = by_id.get(cid)
        if payload is None:
            log.append(f"[map] '{cid}' not produced → will handle per place_unmapped={place_unmapped}")
            continue
        slide_found, shp = _find_shape_anywhere(prs, shape_name)
        if slide_found and shp:
            placed = False
            try:
                if hasattr(shp, "chart") and shp.chart is not None:
                    shp.chart.replace_data(_safe_chart_data(payload))
                    placed = True
            except Exception as e:
                log.append(f"[native:replace_data:error] {type(e).__name__}: {e}")
                placed = False
            if not placed:
                try:
                    if getattr(shp, "is_placeholder", False) and hasattr(shp, "insert_chart"):
                        gf = shp.insert_chart(_xl_chart_type(payload["type"]), _safe_chart_data(payload))
                        placed = gf is not None and getattr(gf, "chart", None) is not None
                except Exception as e:
                    log.append(f"[native:insert_chart:error] {type(e).__name__}: {e}")
                    placed = False
            if not placed:
                try:
                    placed = add_native_chart_to_slide(slide_found, payload, shp.left, shp.top, shp.width, shp.height)
                except Exception as e:
                    log.append(f"[native:add_at_bounds:error] {type(e).__name__}: {e}")
                    placed = False
            if not placed:
                try:
                    slide_found.shapes.add_picture(str(payload["img_path"]), shp.left, shp.top, width=shp.width)
                    log.append(f"[map:fallback-image] '{cid}' → '{shape_name}'")
                    placed = True
                except Exception as e:
                    log.append(f"[map:failed] '{cid}' → '{shape_name}' ({type(e).__name__}: {e})")
            if placed:
                used.add(cid)
        else:
            log.append(f"[map] shape '{shape_name}' not found in template → will handle per place_unmapped={place_unmapped}")

    leftovers = [c for c in produced if c["id"] not in used]
    if place_unmapped and leftovers:
        for c in leftovers:
            slide = prs.slides.add_slide(prs.slide_layouts[title_layout_idx])
            try:
                slide.shapes.title.text = str(c["title"])
            except Exception:
                pass
            ok = add_native_chart_to_slide(slide, c, Inches(1), Inches(1.6), Inches(8), Inches(4.5))
            if not ok:
                slide.shapes.add_picture(str(c["img_path"]), Inches(1), Inches(1.6), width=Inches(8))
                log.append(f"[auto:fallback-image] '{c['id']}'")
    elif not place_unmapped and leftovers:
        log.append(f"[auto:skip-unmapped] skipped {len(leftovers)} leftover chart(s) because mapping exists")

    tgt_slide, tgt_shape = _find_shape_anywhere(prs, insights_target)
    if tgt_shape is not None and hasattr(tgt_shape, "text_frame"):
        tf = tgt_shape.text_frame
        tf.clear()
        for i, bullet in enumerate(s.get("insights", []) or []):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = bullet
            p.level = 0
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[title_layout_idx])
        try:
            slide.shapes.title.text = "Key Insights"
        except Exception:
            pass
        tf = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(4.5)).text_frame
        for i, bullet in enumerate(s.get("insights", []) or []):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = bullet
            p.level = 0
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT

    base, ext = os.path.splitext(os.path.basename(s["ppt_template_path"]))
    out_path = os.path.join(os.path.dirname(s["ppt_template_path"]), f"{base}_filled{ext}")
    prs.save(out_path)
    log.append(f"[ppt] wrote {out_path}")
    return {"ppt_out_path": out_path, "log": log}

def build_app():
    g = StateGraph(State)
    g.add_node("validate", validate_inputs)
    g.add_node("load_mapping", load_mapping)
    g.add_node("load_excel", load_excel)
    g.add_node("plan_charts", plan_charts_with_gemini)
    g.add_node("enforce_requirements", enforce_requirements_overlay)
    g.add_node("render", render_charts_prepare_payloads)
    g.add_node("insights", generate_insights_gemini)
    g.add_node("bind", bind_ppt_native_with_mapping_and_fallback)
    g.set_entry_point("validate")
    g.add_edge("validate", "load_mapping")
    g.add_edge("load_mapping", "load_excel")
    g.add_edge("load_excel", "plan_charts")
    g.add_edge("plan_charts", "enforce_requirements")
    g.add_edge("enforce_requirements", "render")
    g.add_edge("render", "insights")
    g.add_edge("insights", "bind")
    g.add_edge("bind", END)
    return g.compile()

def main():
    ap = argparse.ArgumentParser(description="Excel → Charts (LLM + optional requirements) → Native PPT (shape mapping) + Insights")
    ap.add_argument("--excel", required=True, help="Path to Excel file")
    ap.add_argument("--template", required=True, help="Path to PPTX template")
    ap.add_argument("--map", dest="mapping_path", help="Path to mapping file (.yaml/.yml/.json)")
    args = ap.parse_args()
    app = build_app()
    result: State = app.invoke({
        "excel_path": args.excel,
        "ppt_template_path": args.template,
        "mapping_path": args.mapping_path,
    })
    print("Done!")
    print("Output:", result.get("ppt_out_path"))
    print("Log:")
    for line in result.get("log", []):
        print(" ", line)

if __name__ == "__main__":
    main()
